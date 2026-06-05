"""
Table-specific operations for PowerPoint presentations.
Handles table creation, filling, and cell formatting.
"""

from pptx.util import Pt
from pptx.dml.color import RGBColor
from typing import Dict, Any, Optional, List
import copy


def extract_table_headers(table) -> List[str]:
    """
    Extracts header names from the first row of a table.
    Returns a list of header names.
    """
    if len(table.rows) < 1:
        return []

    headers = []
    header_row = table.rows[0]

    for cell in header_row.cells:
        cell_text = cell.text.strip()
        headers.append(cell_text)

    return headers


def extract_template_row_formatting(table) -> Dict:
    """
    Extracts formatting information from the template row (second row, index 1).
    Returns a dict mapping column index to formatting info.
    """
    if len(table.rows) < 2:
        return {}

    template_row = table.rows[1]
    formatting = {}

    for col_idx, cell in enumerate(template_row.cells):
        if cell.text_frame.paragraphs:
            para = cell.text_frame.paragraphs[0]
            if para.runs:
                run = para.runs[0]
                formatting[col_idx] = {
                    "font_name": run.font.name,
                    "font_size": run.font.size,
                    "font_bold": run.font.bold,
                    "font_italic": run.font.italic,
                    "font_underline": run.font.underline,
                    "alignment": para.alignment
                }

    return formatting


def fill_table(
    table,
    data,
    column_headers: Optional[List[str]] = None,
    field_placeholder: Optional[str] = None,
    template_metadata: Optional[Dict[str, Any]] = None
):
    """
    Fills a PowerPoint table with data, adding rows as needed and preserving styles.

    Supports two data formats:
    1. List of dicts: [{"Name": "John", "Email": "john@example.com"}, ...]
    2. List of lists: [["John", "john@example.com"], ...]

    Assumes row 0 is header and row 1 is template row.

    Args:
        table: The PowerPoint table to fill
        data: List of row data (dicts or lists)
        column_headers: Optional list of column header names
        field_placeholder: Optional field name for marker-based formatting
        template_metadata: Optional template metadata with formatting conventions
    """
    if not data:
        return

    header_rows = 1  # Header is at row 0, template is at row 1

    # Determine data format and convert if needed
    if data and isinstance(data[0], dict):
        rows_to_fill = data
        is_object_format = True
    elif data and isinstance(data[0], list):
        rows_to_fill = data
        is_object_format = False
    else:
        return

    # Fill each row
    for row_idx, row_data in enumerate(rows_to_fill):
        if row_idx == 0:
            current_row = table.rows[header_rows]
        else:
            current_row = add_row_to_table(table)

        # Fill cells in the row
        for col_idx in range(len(table.columns)):
            if col_idx >= len(table.columns):
                break

            cell = current_row.cells[col_idx]

            # Get the cell value based on data format
            if is_object_format:
                if column_headers and col_idx < len(column_headers):
                    header_name = column_headers[col_idx]
                    cell_value = row_data.get(header_name, "")
                else:
                    cell_value = ""
            else:
                if col_idx < len(row_data):
                    cell_value = row_data[col_idx]
                else:
                    cell_value = ""

            # Set cell text while preserving formatting
            set_cell_text_preserve_formatting(
                cell,
                cell_value,
                field_placeholder=field_placeholder,
                template_metadata=template_metadata
            )


def count_template_body_rows(table) -> int:
    """
    Counts the body rows in a pristine template table (all rows minus the
    header row at index 0). This represents how many data rows the template
    designer laid out to fit on the slide.
    """
    return max(0, len(table.rows) - 1)


def reset_table_to_template(table) -> None:
    """
    Resets a table back to its pristine template shape: header row (index 0)
    plus a single empty template row (index 1). Any additional body rows are
    removed and the template row's text is cleared.

    Used on duplicated continuation slides so each can be filled with its own
    chunk of rows while preserving the template's styling.
    """
    tbl = table._tbl
    rows = tbl.tr_lst

    # Keep header (0) and the first template row (1); drop the rest.
    for row_xml in rows[2:]:
        tbl.remove(row_xml)

    # Clear text in the remaining template row's cells.
    if len(tbl.tr_lst) >= 2:
        template_row = tbl.tr_lst[1]
        for tc in template_row.tc_lst:
            for t in tc.xpath('.//a:t'):
                t.text = ""


def add_row_to_table(table):
    """
    Manually adds a row to a table by manipulating the underlying XML.
    python-pptx doesn't have a native table.rows.add() method.
    """
    tbl = table._tbl
    # Create a deep copy of the last row
    last_row_xml = tbl.tr_lst[-1]
    new_row_xml = copy.deepcopy(last_row_xml)

    # Clear text in the new row's cells
    for tc in new_row_xml.tc_lst:
        for t in tc.xpath('.//a:t'):
            t.text = ""

    # Append the new row XML to the table
    tbl.append(new_row_xml)

    # Return the newly created row object
    return table.rows[len(table.rows) - 1]


def set_cell_text_preserve_formatting(
    cell,
    text,
    field_placeholder: Optional[str] = None,
    template_metadata: Optional[Dict[str, Any]] = None
):
    """
    Sets cell text while preserving all existing formatting.
    Handles both strings and arrays of strings (with line breaks between items).
    Applies marker-based formatting if template metadata is provided.

    Args:
        cell: The table cell to update
        text: Text or array of text to insert
        field_placeholder: Optional field name for marker-based formatting
        template_metadata: Optional template metadata with formatting conventions

    IMPORTANT: Do NOT use cell.text = ... because it clears formatting.
    Instead, replace text in the existing runs while keeping their properties.
    """
    if not cell.text_frame.paragraphs:
        return

    text_frame = cell.text_frame

    # Convert text to list if it's a string
    if isinstance(text, str):
        paragraphs_to_add = [text]
    elif isinstance(text, list):
        # Filter out empty strings from array
        paragraphs_to_add = [str(item) for item in text if item]
    else:
        paragraphs_to_add = [str(text)]

    if not paragraphs_to_add:
        # If no content, clear the cell
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.text = ""
        return

    # Create parser for marker-based formatting if metadata provided
    parser = None
    if field_placeholder and template_metadata:
        from app.utils.marker_parser import MarkerParser
        parser = MarkerParser.for_field(template_metadata, field_placeholder)

    # Extract base font properties and paragraph formatting from the first paragraph
    first_paragraph = text_frame.paragraphs[0]
    base_font_props = {}
    base_paragraph_props = {
        "alignment": first_paragraph.alignment,
        "level": first_paragraph.level,
        "space_before": first_paragraph.space_before,
        "space_after": first_paragraph.space_after,
        "line_spacing": first_paragraph.line_spacing
    }

    if first_paragraph.runs:
        first_run = first_paragraph.runs[0]
        base_font_props = {
            "name": first_run.font.name,
            "size": first_run.font.size,
            "bold": first_run.font.bold,
            "italic": first_run.font.italic,
            "underline": first_run.font.underline,
            "color": first_run.font.color.rgb if hasattr(first_run.font.color, 'rgb') else None
        }

    # Set text to the first paragraph with formatting
    if parser:
        # Clear existing runs first
        for run in first_paragraph.runs:
            run.text = ""
        # Import here to avoid circular imports
        from app.services.pptx_service import apply_formatted_text
        apply_formatted_text(first_paragraph, paragraphs_to_add[0], parser, base_font_props)
    else:
        # Plain text without marker formatting
        if first_paragraph.runs:
            first_paragraph.runs[0].text = paragraphs_to_add[0]
            for run in first_paragraph.runs[1:]:
                run.text = ""
        else:
            run = first_paragraph.add_run()
            run.text = paragraphs_to_add[0]

    # Define consistent spacing for all paragraphs in the array
    default_spacing = Pt(18)

    # If we have multiple paragraphs (array), add spacing after first paragraph for visibility
    if len(paragraphs_to_add) > 1:
        if base_paragraph_props["space_after"] is not None and base_paragraph_props["space_after"] > Pt(0):
            first_paragraph.space_after = base_paragraph_props["space_after"]
        else:
            first_paragraph.space_after = default_spacing

    # Add remaining paragraphs for array items
    for idx, paragraph_text in enumerate(paragraphs_to_add[1:]):
        # Add a new paragraph
        new_paragraph = text_frame.add_paragraph()

        # Apply paragraph formatting
        if base_paragraph_props["alignment"] is not None:
            new_paragraph.alignment = base_paragraph_props["alignment"]
        if base_paragraph_props["level"] is not None:
            new_paragraph.level = base_paragraph_props["level"]
        if base_paragraph_props["space_before"] is not None:
            new_paragraph.space_before = base_paragraph_props["space_before"]

        # Add explicit spacing after paragraphs for visibility
        if base_paragraph_props["space_after"] is not None and base_paragraph_props["space_after"] > Pt(0):
            new_paragraph.space_after = base_paragraph_props["space_after"]
        else:
            new_paragraph.space_after = default_spacing

        if base_paragraph_props["line_spacing"] is not None:
            new_paragraph.line_spacing = base_paragraph_props["line_spacing"]

        if parser:
            # Import here to avoid circular imports
            from app.services.pptx_service import apply_formatted_text
            apply_formatted_text(new_paragraph, paragraph_text, parser, base_font_props)
        else:
            # Plain text - preserve formatting from the first paragraph
            if first_paragraph.runs:
                first_run = first_paragraph.runs[0]
                # Add run with text and copy ALL formatting from the template
                run = new_paragraph.add_run(paragraph_text)
                # Use copy_run_formatting to properly handle all font properties
                from app.services.pptx_service import copy_run_formatting
                copy_run_formatting(first_run, run)
            else:
                # Fallback if no runs in template
                new_paragraph.text = paragraph_text


def apply_cell_formatting(cell, formatting):
    """
    Applies formatting to a cell based on a formatting dict.
    """
    if not cell.text_frame.paragraphs:
        return

    para = cell.text_frame.paragraphs[0]

    if para.runs:
        run = para.runs[0]
        font = run.font

        if formatting.get("font_name"):
            font.name = formatting["font_name"]
        if formatting.get("font_size"):
            font.size = formatting["font_size"]
        if formatting.get("font_bold") is not None:
            font.bold = formatting["font_bold"]
        if formatting.get("font_italic") is not None:
            font.italic = formatting["font_italic"]
        if formatting.get("font_underline") is not None:
            font.underline = formatting["font_underline"]

    if formatting.get("alignment") is not None:
        para.alignment = formatting["alignment"]


def copy_cell_style(source_cell, target_cell):
    """
    Copies font styling from one cell to another.
    """
    if not source_cell.text_frame.paragraphs:
        return

    source_para = source_cell.text_frame.paragraphs[0]
    target_para = target_cell.text_frame.paragraphs[0]

    # Copy paragraph alignment
    target_para.alignment = source_para.alignment

    if source_para.runs and target_para.runs:
        source_run = source_para.runs[0]
        target_run = target_para.runs[0]

        target_run.font.name = source_run.font.name
        target_run.font.size = source_run.font.size
        target_run.font.bold = source_run.font.bold
        target_run.font.italic = source_run.font.italic
        if source_run.font.color and source_run.font.color.type != 0:
            try:
                target_run.font.color.rgb = source_run.font.color.rgb
            except:
                pass
