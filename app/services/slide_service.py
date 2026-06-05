"""
Slide-level structural operations for PowerPoint presentations.

python-pptx has no native slide duplication, so this module deep-copies a
slide's shape tree into a new slide and rewires its relationships (images,
etc.) so the copy renders identically to the source.
"""

import copy
import logging
from typing import Optional

from pptx.oxml.ns import qn
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

# Relationship namespace used by attributes like r:embed / r:link / r:id
_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

# Relationships that the new slide already owns (created by add_slide) and must
# therefore NOT be copied from the source slide.
_SKIP_RELTYPES = {RT.SLIDE_LAYOUT}


def duplicate_slide_after(presentation, source_slide, after_slide):
    """
    Creates a faithful copy of ``source_slide`` and positions it immediately
    after ``after_slide`` in the slide order.

    The copy preserves every shape (text, tables, pictures) along with its
    formatting, and rewires image/media relationships so they keep working.

    Args:
        presentation: The Presentation object.
        source_slide: The slide whose content should be copied.
        after_slide: The slide that the new copy should follow.

    Returns:
        The newly created slide.
    """
    dest_slide = presentation.slides.add_slide(source_slide.slide_layout)

    # add_slide() seeds the slide with placeholders inherited from the layout.
    # Remove them so we start from a clean shape tree before copying.
    for shape in list(dest_slide.shapes):
        shape._element.getparent().remove(shape._element)

    # Copy the source slide's relationships into the destination part, building
    # a map of old rId -> new rId. python-pptx assigns new rIds on add, so the
    # copied shape XML must be remapped afterwards.
    rId_map = _copy_relationships(source_slide, dest_slide)

    # Deep-copy every shape element from the source into the destination.
    dest_spTree = dest_slide.shapes._spTree
    for shape in source_slide.shapes:
        new_element = copy.deepcopy(shape._element)
        dest_spTree.insert_element_before(new_element, "p:extLst")

    # Point the copied shapes at the destination part's relationships.
    _remap_relationship_ids(dest_spTree, rId_map)

    # Reorder so the duplicate sits right after the anchor slide.
    _move_slide_after(presentation, after_slide, dest_slide)

    return dest_slide


def _copy_relationships(source_slide, dest_slide) -> dict:
    """
    Copies source slide relationships onto the destination slide part.
    Returns a mapping of {old_rId: new_rId}.
    """
    rId_map = {}
    for rId, rel in source_slide.part.rels.items():
        if rel.reltype in _SKIP_RELTYPES:
            continue
        try:
            if rel.is_external:
                new_rId = dest_slide.part.rels._add_relationship(
                    rel.reltype, rel.target_ref, is_external=True
                )
            else:
                new_rId = dest_slide.part.rels._add_relationship(
                    rel.reltype, rel.target_part
                )
            rId_map[rId] = new_rId
        except Exception as exc:  # pragma: no cover - defensive
            logging.error(f"Failed to copy relationship {rId} ({rel.reltype}): {exc}")
    return rId_map


def _remap_relationship_ids(spTree, rId_map: dict) -> None:
    """
    Walks the copied shape tree and rewrites any relationship-id attribute
    (r:embed, r:link, r:id, ...) to its new value in the destination part.
    """
    if not rId_map:
        return

    prefix = "{%s}" % _R_NS
    for element in spTree.iter():
        for attr_name, attr_value in list(element.attrib.items()):
            if attr_name.startswith(prefix) and attr_value in rId_map:
                element.set(attr_name, rId_map[attr_value])


def _move_slide_after(presentation, after_slide, slide_to_move) -> None:
    """
    Moves ``slide_to_move`` so it directly follows ``after_slide`` in the
    presentation's slide ordering.
    """
    sldIdLst = presentation.slides._sldIdLst

    after_element = _sldId_element_for(presentation, after_slide)
    move_element = _sldId_element_for(presentation, slide_to_move)

    if after_element is None or move_element is None:
        logging.warning("Could not reorder duplicated slide; leaving at end.")
        return

    sldIdLst.remove(move_element)
    after_element.addnext(move_element)


def _sldId_element_for(presentation, slide):
    """
    Finds the <p:sldId> element in the presentation that references ``slide``.
    """
    for sldId in presentation.slides._sldIdLst:
        rId = sldId.get(qn("r:id"))
        try:
            if presentation.part.rels[rId].target_part is slide.part:
                return sldId
        except KeyError:
            continue
    return None
