from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import re
import shutil
import subprocess
from pathlib import Path
from typing import Dict, Any, Optional
from app.utils.placeholder import infer_placeholder_type, extract_placeholder_metadata
from app.utils.map_logic import apply_map_logic
from app.utils.pptx_utils import get_shape_alt_text
from app.services.image_service import get_company_logo_path, get_topic_image_path, cleanup_temp_images
from app.services.table_service import (
    fill_table,
    extract_table_headers,
    count_template_body_rows,
    reset_table_to_template,
)
from app.services.slide_service import duplicate_slide_after
import logging

# metadata extraction and presentation generation logic will be use by an ai agent to create pptx files based on user input and a template file. The pptx template will have placeholders like {{title}}, {{summary}}, etc. which will be replaced by the ai agent with actual content before generating the final presentation.
def extract_ppt_metadata(template_path: str):

    presentation = Presentation(template_path)

    slides_data = []

    pattern = r"\{\{(.*?)\}\}"

    for slide_index, slide in enumerate(presentation.slides):

        slide_info = {
            "slide_number": slide_index + 1,
            "placeholders": []
        }
        # logging.info(f"Slide {slide_index + 1}")

        for shape_index, shape in enumerate(slide.shapes):
            # logging.info(f"Shape {shape_index} in slide {slide_index + 1}")
            # logging.info(f"Shape: {shape}")
            # logging.info(f"Shape text: {getattr(shape, 'text', 'No text attribute')}")
            # logging.info(f"Shape has text: {hasattr(shape, 'text')}")

            # 1. Check text for placeholders
            shape_text = getattr(shape, "text", None)
            if shape_text and shape_text.strip():
                matches = re.findall(pattern, shape_text.strip())
                for match in matches:
                    metadata = extract_placeholder_metadata(match)
                    placeholder_name = metadata["name"]
                    paragraphs_count = metadata["paragraphs"]
                    placeholder_type = infer_placeholder_type(placeholder_name)
                    slide_info["placeholders"].append({
                        "placeholder": placeholder_name,
                        "slide_number": slide_index + 1,
                        "shape_index": shape_index,
                        "type": placeholder_type,
                        "paragraphs": paragraphs_count
                    })

            # 2. Check Alt Text for placeholders (common for images)
            alt_text = get_shape_alt_text(shape)
            if alt_text:
                matches = re.findall(pattern, alt_text)
                for match in matches:
                    metadata = extract_placeholder_metadata(match)
                    placeholder_name = metadata["name"]
                    paragraphs_count = metadata["paragraphs"]

                    # Avoid duplicate detection if it was already in the text (rare)
                    if any(p["placeholder"] == placeholder_name and p["shape_index"] == shape_index for p in slide_info["placeholders"]):
                        continue

                    placeholder_type = infer_placeholder_type(placeholder_name)
                    slide_info["placeholders"].append({
                        "placeholder": placeholder_name,
                        "slide_number": slide_index + 1,
                        "shape_index": shape_index,
                        "type": placeholder_type,
                        "paragraphs": paragraphs_count
                    })

            # 3. Check for Tables with placeholders in Alt Text
            if shape.has_table:
                alt_text = get_shape_alt_text(shape)
                if alt_text:
                    matches = re.findall(pattern, alt_text)
                    for match in matches:
                        metadata = extract_placeholder_metadata(match)
                        placeholder_name = metadata["name"]
                        placeholder_type = infer_placeholder_type(placeholder_name)
                        if placeholder_type == "table":
                            # Extract header names from the first row
                            table = getattr(shape, "table", None)
                            if not table:
                                continue
                                
                            column_headers = extract_table_headers(table)

                            slide_info["placeholders"].append({
                                "placeholder": placeholder_name,
                                "slide_number": slide_index + 1,
                                "shape_index": shape_index,
                                "type": "table",
                                "columns": len(table.columns),
                                "column_headers": column_headers,
                                "max_chars": 500,
                                "paragraphs": 1
                            })


        slides_data.append(slide_info)

    return slides_data

def insert_paragraph_after(reference_paragraph, text_frame):
    """
    Inserts a new, empty paragraph immediately AFTER the given reference paragraph
    in the text frame, and returns it as a python-pptx paragraph object.

    python-pptx's text_frame.add_paragraph() always appends to the END of the text
    frame. When a template has a trailing empty paragraph after the placeholder line,
    appending puts new paragraphs after that empty one, stranding it between the
    first and second generated paragraphs and creating an inconsistent gap.
    Inserting right after the reference paragraph keeps generated paragraphs contiguous.
    """
    from pptx.oxml.ns import qn
    import copy as _copy

    ref_p = reference_paragraph._p
    # Build a fresh empty <a:p> element (no runs, no pPr) to avoid inheriting
    # run-level content; paragraph-level props are applied by the caller.
    new_p = ref_p.makeelement(qn('a:p'), {})
    ref_p.addnext(new_p)

    # Wrap the raw XML element in a python-pptx _Paragraph proxy.
    from pptx.text.text import _Paragraph
    return _Paragraph(new_p, reference_paragraph._parent)


def remove_empty_paragraphs(text_frame):
    """
    Removes empty paragraphs (no text content) from a text frame.
    Used to clean up leftover template paragraphs that would otherwise create
    inconsistent vertical spacing between generated paragraphs.
    """
    paragraphs = text_frame.paragraphs
    # Never remove the very last paragraph if it's the only one left.
    for p in list(paragraphs):
        if not p.text.strip() and len(text_frame.paragraphs) > 1:
            p._p.getparent().remove(p._p)


def copy_run_formatting(source_run, target_run):
    """
    Copies all font formatting from source_run to target_run.
    Only copies properties that are explicitly set (not None) to avoid resetting to defaults.
    """
    source_font = source_run.font
    target_font = target_run.font

    # Copy properties only if they are explicitly set (not None/inherited)
    if source_font.name is not None:
        target_font.name = source_font.name
    if source_font.size is not None:
        target_font.size = source_font.size
    if source_font.bold is not None:
        target_font.bold = source_font.bold
    if source_font.italic is not None:
        target_font.italic = source_font.italic
    if source_font.underline is not None:
        target_font.underline = source_font.underline

    # Copy color if it exists
    if source_font.color:
        try:
            if hasattr(source_font.color, 'rgb'):
                target_font.color.rgb = source_font.color.rgb
            elif hasattr(source_font.color, 'type'):
                target_font.color.type = source_font.color.type
        except:
            pass


def apply_formatted_text(paragraph, text: str, parser, base_font_props: Optional[Dict[str, Any]] = None) -> None:
    """
    Apply text with parsed marker formatting to a PowerPoint paragraph.
    Preserves the placeholder's base font and applies marker-based formatting only to marked text.

    Args:
        paragraph: The paragraph to add formatted text to
        text: Text potentially containing markers like $$marker$$content$$marker$$
        parser: MarkerParser instance with field-specific conventions
        base_font_props: Optional base font properties to apply to all runs
    """
    # Extract base font properties from the first existing run if not provided
    if not base_font_props:
        base_font_props = {}
        if paragraph.runs:
            first_run = paragraph.runs[0]
            base_font_props = {
                "name": first_run.font.name,
                "size": first_run.font.size,
                "bold": first_run.font.bold,
                "italic": first_run.font.italic,
                "underline": first_run.font.underline,
                "color": first_run.font.color.rgb if hasattr(first_run.font.color, 'rgb') else None
            }

    runs_data = parser.get_pptx_runs(text)

    # Clear existing runs
    for run in list(paragraph.runs):
        run.text = ""

    for run_data in runs_data:
        run = paragraph.add_run()
        run.text = run_data["text"]

        # Start with base font properties (placeholder font)
        # Only apply if the property was explicitly set (not None)
        if base_font_props.get("name") is not None:
            run.font.name = base_font_props["name"]
        if base_font_props.get("size") is not None:
            run.font.size = base_font_props["size"]
        if base_font_props.get("bold") is not None:
            run.font.bold = base_font_props["bold"]
        if base_font_props.get("italic") is not None:
            run.font.italic = base_font_props["italic"]
        if base_font_props.get("underline") is not None:
            run.font.underline = base_font_props["underline"]
        if base_font_props.get("color") is not None:
            run.font.color.rgb = base_font_props["color"]

        # Override with marker-specific formatting (only if marker style is defined)
        # Color from marker overrides placeholder color
        if run_data.get("color"):
            try:
                r = int(run_data["color"][0:2], 16)
                g = int(run_data["color"][2:4], 16)
                b = int(run_data["color"][4:6], 16)
                run.font.color.rgb = RGBColor(r, g, b)
            except (ValueError, TypeError):
                pass

        # Bold from marker overrides placeholder bold (only if marker specifies bold)
        if run_data.get("bold"):
            run.font.bold = True

        # Italic from marker overrides placeholder italic (only if marker specifies italic)
        if run_data.get("italic"):
            run.font.italic = True

        # Font size from marker overrides placeholder size (only if marker specifies size)
        if run_data.get("font_size"):
            try:
                run.font.size = Pt(int(run_data["font_size"]))
            except (ValueError, TypeError):
                pass


def replace_text_preserve_formatting(
    shape,
    placeholder: str,
    value,
    field_placeholder: Optional[str] = None,
    template_metadata: Optional[Dict[str, Any]] = None
):
    """
    Replaces placeholder text in a shape while preserving the original formatting.
    Optionally applies field-specific formatting conventions based on markers.

    Args:
        shape: The PowerPoint shape containing the placeholder
        placeholder: The placeholder text to find and replace
        value: The replacement value (str or list of str)
        field_placeholder: The field name for marker parsing (e.g., "program_schedule")
        template_metadata: Template metadata with formatting_conventions
    """
    if not hasattr(shape, "text_frame"):
        return False

    text_frame = shape.text_frame

    # If field context provided, use marker-based formatting
    if field_placeholder and template_metadata:
        from app.utils.marker_parser import MarkerParser

        parser = MarkerParser.for_field(template_metadata, field_placeholder)

        # Handle list of values
        if isinstance(value, list):
            # Find which paragraph contains the placeholder
            placeholder_para_idx = None

            for para_idx, paragraph in enumerate(text_frame.paragraphs):
                if placeholder in paragraph.text:
                    placeholder_para_idx = para_idx
                    break

            if placeholder_para_idx is not None:
                template_paragraph = text_frame.paragraphs[placeholder_para_idx]

                # Extract base font properties and paragraph formatting from template
                base_font_props = {}
                base_paragraph_props = {
                    "alignment": template_paragraph.alignment,
                    "level": template_paragraph.level,
                    "space_before": template_paragraph.space_before,
                    "space_after": template_paragraph.space_after,
                    "line_spacing": template_paragraph.line_spacing
                }
                if template_paragraph.runs:
                    first_run = template_paragraph.runs[0]
                    base_font_props = {
                        "name": first_run.font.name,
                        "size": first_run.font.size,
                        "bold": first_run.font.bold,
                        "italic": first_run.font.italic,
                        "underline": first_run.font.underline,
                        "color": first_run.font.color.rgb if hasattr(first_run.font.color, 'rgb') else None
                    }

                # Define consistent spacing for all paragraphs
                from pptx.util import Pt
                default_spacing = Pt(18)

                def _apply_para_props(p):
                    # Apply uniform paragraph-level formatting so every gap is identical.
                    if base_paragraph_props["alignment"] is not None:
                        p.alignment = base_paragraph_props["alignment"]
                    if base_paragraph_props["level"] is not None:
                        p.level = base_paragraph_props["level"]
                    if base_paragraph_props["space_before"] is not None:
                        p.space_before = base_paragraph_props["space_before"]
                    if base_paragraph_props["space_after"] is not None and base_paragraph_props["space_after"] > Pt(0):
                        p.space_after = base_paragraph_props["space_after"]
                    else:
                        p.space_after = default_spacing
                    if base_paragraph_props["line_spacing"] is not None:
                        p.line_spacing = base_paragraph_props["line_spacing"]

                # Add each value with formatting. New paragraphs are inserted
                # immediately after the previous one so they stay contiguous and
                # the template's trailing empty paragraph cannot land between them.
                prev_para = None
                for i, val in enumerate(value):
                    if i == 0:
                        para = template_paragraph
                        if len(value) > 1:
                            _apply_para_props(para)
                    else:
                        para = insert_paragraph_after(prev_para, text_frame)
                        _apply_para_props(para)

                    # Apply formatted text with base font properties
                    apply_formatted_text(para, str(val), parser, base_font_props)
                    prev_para = para

                # Remove any leftover empty template paragraphs (e.g. trailing blank
                # line after the placeholder) that would create uneven spacing.
                remove_empty_paragraphs(text_frame)

                return True
        else:
            # Handle single value with marker formatting - use same logic as regular path
            # but apply formatting from parser
            replaced = False

            for paragraph in text_frame.paragraphs:
                # Build list of (run, start_index, end_index) for text search
                run_info = []
                cumulative_length = 0

                for run in paragraph.runs:
                    start = cumulative_length
                    end = start + len(run.text)
                    run_info.append((run, start, end))
                    cumulative_length = end

                # Get the full paragraph text
                full_text = paragraph.text

                # Find placeholder in full text
                placeholder_start = full_text.find(placeholder)
                if placeholder_start != -1:
                    placeholder_end = placeholder_start + len(placeholder)

                    # Find which runs this placeholder spans
                    affected_runs = []
                    for run, start, end in run_info:
                        if not (end <= placeholder_start or start >= placeholder_end):
                            affected_runs.append((run, start, end))

                    if affected_runs:
                        # Get formatting info from parser
                        parsed_runs_data = parser.get_pptx_runs(str(value))
                        replacement_str = str(value)

                        # If placeholder is within a single run (most common case)
                        if len(affected_runs) == 1:
                            run, run_start, run_end = affected_runs[0]
                            # Calculate position within the run
                            offset_start = placeholder_start - run_start
                            offset_end = offset_start + len(placeholder)
                            # Replace just the placeholder, keep rest of text
                            run.text = run.text[:offset_start] + replacement_str + run.text[offset_end:]

                            # Apply formatting from parsed runs if available
                            if len(parsed_runs_data) > 0:
                                run_data = parsed_runs_data[0]
                                if run_data.get("color"):
                                    try:
                                        r = int(run_data["color"][0:2], 16)
                                        g = int(run_data["color"][2:4], 16)
                                        b = int(run_data["color"][4:6], 16)
                                        run.font.color.rgb = RGBColor(r, g, b)
                                    except (ValueError, TypeError):
                                        pass
                                if run_data.get("bold"):
                                    run.font.bold = True
                                if run_data.get("italic"):
                                    run.font.italic = True
                        else:
                            # Placeholder spans multiple runs - clear all and put replacement in first
                            for run, _, _ in affected_runs:
                                run.text = ""
                            first_run = affected_runs[0][0]
                            first_run.text = replacement_str

                            # Apply formatting
                            if len(parsed_runs_data) > 0:
                                run_data = parsed_runs_data[0]
                                if run_data.get("color"):
                                    try:
                                        r = int(run_data["color"][0:2], 16)
                                        g = int(run_data["color"][2:4], 16)
                                        b = int(run_data["color"][4:6], 16)
                                        first_run.font.color.rgb = RGBColor(r, g, b)
                                    except (ValueError, TypeError):
                                        pass
                                if run_data.get("bold"):
                                    first_run.font.bold = True
                                if run_data.get("italic"):
                                    first_run.font.italic = True

                        replaced = True

            return replaced

    # Handle list of paragraphs (multi-paragraph values)
    if isinstance(value, list) and len(text_frame.paragraphs) > 0:
        # Find which paragraph contains the placeholder
        placeholder_para_idx = None
        template_run = None

        for para_idx, paragraph in enumerate(text_frame.paragraphs):
            if placeholder in paragraph.text:
                placeholder_para_idx = para_idx
                # Get formatting from first run
                if paragraph.runs:
                    template_run = paragraph.runs[0]
                break

        if placeholder_para_idx is not None:
            template_paragraph = text_frame.paragraphs[placeholder_para_idx]

            # Extract paragraph formatting from template
            base_paragraph_props = {
                "alignment": template_paragraph.alignment,
                "level": template_paragraph.level,
                "space_before": template_paragraph.space_before,
                "space_after": template_paragraph.space_after,
                "line_spacing": template_paragraph.line_spacing
            }

            # Define consistent spacing for all paragraphs
            from pptx.util import Pt
            default_spacing = Pt(18)

            def _apply_para_props(p):
                # Apply uniform paragraph-level formatting so every gap is identical.
                if base_paragraph_props["alignment"] is not None:
                    p.alignment = base_paragraph_props["alignment"]
                if base_paragraph_props["level"] is not None:
                    p.level = base_paragraph_props["level"]
                if base_paragraph_props["space_before"] is not None:
                    p.space_before = base_paragraph_props["space_before"]
                if base_paragraph_props["space_after"] is not None and base_paragraph_props["space_after"] > Pt(0):
                    p.space_after = base_paragraph_props["space_after"]
                else:
                    p.space_after = default_spacing
                if base_paragraph_props["line_spacing"] is not None:
                    p.line_spacing = base_paragraph_props["line_spacing"]

            # Clear the placeholder paragraph
            for run in list(template_paragraph.runs):
                run.text = ""

            # Add each paragraph from the list. New paragraphs are inserted
            # immediately after the previous one so they stay contiguous and the
            # template's trailing empty paragraph cannot land between them.
            prev_para = None
            for i, para_text in enumerate(value):
                if i == 0:
                    # Use the existing paragraph with the placeholder
                    para = template_paragraph
                    if len(value) > 1:
                        _apply_para_props(para)
                else:
                    para = insert_paragraph_after(prev_para, text_frame)
                    _apply_para_props(para)

                # Clear and add text
                for run in list(para.runs):
                    run.text = ""

                run = para.add_run()
                run.text = str(para_text)

                # Apply template formatting
                if template_run:
                    copy_run_formatting(template_run, run)

                prev_para = para

            # Remove any leftover empty template paragraphs (e.g. trailing blank
            # line after the placeholder) that would create uneven spacing.
            remove_empty_paragraphs(text_frame)

            return True

        return False

    # Handle single text value (original logic)
    replaced = False

    # Iterate through all paragraphs and runs
    for paragraph in text_frame.paragraphs:
        # Build a list of (run, start_index, end_index) for text search
        run_info = []
        cumulative_length = 0

        for run in paragraph.runs:
            start = cumulative_length
            end = start + len(run.text)
            run_info.append((run, start, end))
            cumulative_length = end

        # Get the full paragraph text
        full_text = paragraph.text

        # Find placeholder in full text
        placeholder_start = full_text.find(placeholder)
        if placeholder_start != -1:
            placeholder_end = placeholder_start + len(placeholder)

            # Find which runs this placeholder spans
            affected_runs = []
            for run, start, end in run_info:
                if not (end <= placeholder_start or start >= placeholder_end):
                    affected_runs.append((run, start, end))

            if affected_runs:
                # Get the first run's formatting as reference
                reference_run = affected_runs[0][0]
                replacement_str = str(value)

                # If placeholder is within a single run (most common case)
                if len(affected_runs) == 1:
                    run, run_start, run_end = affected_runs[0]
                    # Calculate position within the run
                    offset_start = placeholder_start - run_start
                    offset_end = offset_start + len(placeholder)
                    # Replace just the placeholder, keep rest of text
                    run.text = run.text[:offset_start] + replacement_str + run.text[offset_end:]
                    copy_run_formatting(reference_run, run)
                else:
                    # Placeholder spans multiple runs - clear all and put replacement in first
                    for run, _, _ in affected_runs:
                        run.text = ""
                    first_run = affected_runs[0][0]
                    first_run.text = replacement_str
                    copy_run_formatting(reference_run, first_run)

                replaced = True

    return replaced

def _determine_rows_per_slide(table, template_metadata, table_name):
    """
    Decides how many data rows fit on a single slide for a given table.

    Resolution order:
      1. Explicit `rows_per_slide` in the table's `field_instructions` entry,
         keyed by the placeholder name (e.g. "table:programme"). This is the
         convention for configuring pagination.
      2. Inferred from the template table's own body-row count (the capacity
         the designer drew). Falls back to 1.
    """
    metadata = template_metadata or {}

    # 1. field_instructions configuration (keyed by placeholder name).
    entry = metadata.get("field_instructions", {}).get(table_name, {})
    if isinstance(entry, dict):
        rps = entry.get("rows_per_slide")
        if isinstance(rps, int) and rps > 0:
            return rps

    # 2. Infer from the template's pristine body-row count.
    return max(1, count_template_body_rows(table))


def _find_table_shape(slide, table_name):
    """
    Finds the table shape on a slide whose placeholder matches table_name.
    """
    pattern = r"\{\{(.*?)\}\}"
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        alt_text = get_shape_alt_text(shape)
        if not alt_text:
            continue
        for match in re.findall(pattern, alt_text):
            if extract_placeholder_metadata(match)["name"] == table_name:
                return shape
    return None


def _fill_tables_with_pagination(presentation, table_jobs, template_metadata):
    """
    Fills each collected table job, splitting rows across duplicated slides
    when they exceed the per-slide capacity. Continuation slides are exact
    copies of the source slide (shared content preserved), with their table
    reset to the template shape and filled with the next chunk of rows.
    """
    for job in table_jobs:
        data = job["table_data"]
        rows_per_slide = job["rows_per_slide"]
        source_slide = job["slide"]
        table_name = job["table_name"]
        column_headers = job["column_headers"]

        # Split rows into per-slide chunks.
        chunks = [
            data[i:i + rows_per_slide]
            for i in range(0, len(data), rows_per_slide)
        ]
        if not chunks:
            continue

        # Duplicate the (still pristine) source slide once per extra chunk,
        # before filling the original, so every copy starts clean.
        extra_slides = []
        anchor = source_slide
        for _ in chunks[1:]:
            try:
                duplicate = duplicate_slide_after(presentation, source_slide, anchor)
                extra_slides.append(duplicate)
                anchor = duplicate
            except Exception as exc:
                logging.error(f"Failed to duplicate slide for table '{table_name}': {exc}")
                extra_slides.append(None)

        # Fill the original slide's table with the first chunk. Reset first so
        # the row count matches the chunk exactly (no leftover template rows).
        reset_table_to_template(job["shape"].table)
        fill_table(
            job["shape"].table,
            chunks[0],
            column_headers=column_headers,
            field_placeholder=table_name,
            template_metadata=template_metadata,
        )

        # Fill each continuation slide with its chunk.
        for duplicate, chunk in zip(extra_slides, chunks[1:]):
            if duplicate is None:
                continue
            dup_shape = _find_table_shape(duplicate, table_name)
            if dup_shape is None:
                logging.warning(
                    f"Continuation slide missing table '{table_name}'; skipping chunk."
                )
                continue
            reset_table_to_template(dup_shape.table)
            fill_table(
                dup_shape.table,
                chunk,
                column_headers=column_headers,
                field_placeholder=table_name,
                template_metadata=template_metadata,
            )


def generate_presentation(
    template_path: str,
    output_path: str,
    replacements: dict,
    template_metadata: Optional[Dict[str, Any]] = None
):

    presentation = Presentation(template_path)
    pattern = r"\{\{(.*?)\}\}"
    shapes_to_remove = []
    # Tables are filled in a separate pass AFTER shape cleanup so that
    # continuation slides (when rows overflow one slide) are duplicated from a
    # clean, fully-rendered slide. Each job captures everything needed to fill.
    table_jobs = []

    # Snapshot the slide list: the pagination pass adds slides, and we must not
    # re-process those generated slides in this loop.
    for slide in list(presentation.slides):
        for shape in slide.shapes:
            # 1. Check Alt Text for Image Placeholders
            alt_text = get_shape_alt_text(shape)
            if alt_text:
                matches = re.findall(pattern, alt_text)
                for match in matches:
                    if "image:logo" in match.lower():
                        domain = replacements.get(match)
                        if domain:
                            logo_path = get_company_logo_path(domain)
                            if logo_path:
                                # Add the picture at the same position and size
                                slide.shapes.add_picture(
                                    logo_path,
                                    shape.left,
                                    shape.top,
                                    width=shape.width,
                                    height=shape.height
                                )
                                shapes_to_remove.append((slide, shape))
                    elif "image:topic" in match.lower() or "image:bg" in match.lower():
                        query = replacements.get(match)
                        if query:
                            img_path = get_topic_image_path(query)
                            if img_path:
                                # Add the picture at the same position and size, covering exactly the shape area
                                slide.shapes.add_picture(
                                    img_path,
                                    shape.left,
                                    shape.top,
                                    width=shape.width,
                                    height=shape.height
                                )
                                shapes_to_remove.append((slide, shape))

            # 2. Check Text for Text Placeholders
            shape_text = getattr(shape, "text", None)
            if shape_text:
                for key, value in replacements.items():
                    # Look for placeholder with or without metadata
                    # e.g., {{key}} or {{key:paragraphs=2}}
                    placeholder_pattern = r"\{\{" + re.escape(key) + r"(?::paragraphs=\d+)?\}\}"
                    match = re.search(placeholder_pattern, shape_text)

                    if match:
                        actual_placeholder = match.group(0)
                        # Extract field placeholder name (without :metadata)
                        field_placeholder = key.split(":")[0]

                        # Handle different value types
                        if isinstance(value, list):
                            # If it's a table (list of lists), skip - handled separately
                            if value and isinstance(value[0], list):
                                continue
                            # For paragraph arrays, pass the list directly
                            # For bullet lists (list of strings), join them
                            # The function will handle both cases
                            replace_text_preserve_formatting(
                                shape,
                                actual_placeholder,
                                value,
                                field_placeholder=field_placeholder,
                                template_metadata=template_metadata
                            )
                        else:
                            replace_text_preserve_formatting(
                                shape,
                                actual_placeholder,
                                str(value),
                                field_placeholder=field_placeholder,
                                template_metadata=template_metadata
                            )

            
            # 3. Check for Tables
            if shape.has_table:
                alt_text = get_shape_alt_text(shape)
                if alt_text:
                    matches = re.findall(pattern, alt_text)
                    for match in matches:
                        if "table:" in match.lower():
                            # Extract just the table name (without metadata)
                            metadata = extract_placeholder_metadata(match)
                            table_name = metadata["name"]

                            table_data = replacements.get(table_name)
                            if table_data and isinstance(table_data, list):
                                # Defer the actual fill until after cleanup so
                                # overflow rows can paginate onto duplicated
                                # slides. Capture values from the pristine table.
                                table = getattr(shape, "table", None)
                                if not table:
                                    continue
                                table_jobs.append({
                                    "slide": slide,
                                    "shape": shape,
                                    "table_name": table_name,
                                    "table_data": table_data,
                                    "column_headers": extract_table_headers(table),
                                    "rows_per_slide": _determine_rows_per_slide(
                                        table, template_metadata, table_name
                                    ),
                                })

        # Apply custom map highlighting and pointer logic
        apply_map_logic(slide, replacements)

    # Clean up: Remove original placeholder shapes that were replaced by images
    for slide, shape in shapes_to_remove:
        try:
            sp = shape._element
            sp.getparent().remove(sp)
        except Exception as e:
            logging.error(f"Error removing shape: {e}")

    # Fill tables, paginating onto duplicated slides when rows overflow.
    _fill_tables_with_pagination(presentation, table_jobs, template_metadata)

    presentation.save(output_path)
    cleanup_temp_images()
    return output_path


def attach_placeholder_values(metadata: list[dict], replacements: dict):

    return [
        {
            **slide,
            "placeholders": [
                {
                    **placeholder,
                    "value": replacements.get(placeholder["placeholder"]) or ""
                }
                for placeholder in slide["placeholders"]
            ]
        }
        for slide in metadata
    ]

def _find_soffice_path():
    soffice_path = shutil.which("soffice") or shutil.which("libreoffice")

    if soffice_path:
        return soffice_path

    possible_paths = [
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ]

    for possible_path in possible_paths:
        if Path(possible_path).is_file():
            return possible_path

    return None

def convert_pptx_to_pdf(pptx_path: str, output_dir: str = "generated"):
    soffice_path = _find_soffice_path()

    if not soffice_path:
        raise RuntimeError(
            "LibreOffice is required to convert PPTX files to PDF. "
            "Install LibreOffice or add soffice to PATH."
        )

    pptx_file = Path(pptx_path).resolve()
    pdf_dir = Path(output_dir).resolve()
    pdf_dir.mkdir(parents=True, exist_ok=True)

    result = subprocess.run(
        [
            soffice_path,
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(pdf_dir),
            str(pptx_file),
        ],
        capture_output=True,
        text=True,
        timeout=60,
        check=False
    )

    pdf_path = pdf_dir / f"{pptx_file.stem}.pdf"

    if result.returncode != 0 or not pdf_path.is_file():
        raise RuntimeError(
            f"PPTX to PDF conversion failed: {result.stderr or result.stdout}"
        )

    return str(pdf_path)

def generate_template_thumbnail(template_path: str, output_dir: str = "public/thumbnails"):
    """
    Generates a PNG thumbnail of the first slide of a PPTX template.
    """
    soffice_path = _find_soffice_path()
    if not soffice_path:
        return None

    template_file = Path(template_path).resolve()
    out_dir = Path(output_dir).resolve()
    out_dir.mkdir(parents=True, exist_ok=True)
    
    thumb_path = out_dir / f"{template_file.stem}.png"
    
    # If thumbnail already exists, don't regenerate (unless we want to force)
    if thumb_path.is_file():
        return str(thumb_path)

    # LibreOffice conversion to PNG (exports slides as images)
    subprocess.run(
        [
            soffice_path,
            "--headless",
            "--convert-to",
            "png",
            "--outdir",
            str(out_dir),
            str(template_file),
        ],
        capture_output=True,
        timeout=30
    )
    
    # LibreOffice might name it template_name-0.png or template_name.png
    # Usually it's template_name.png for single slide or if it just does the first
    possible_names = [f"{template_file.stem}.png", f"{template_file.stem}-0.png"]
    for name in possible_names:
        found = out_dir / name
        if found.is_file():
            if name != f"{template_file.stem}.png":
                found.rename(thumb_path)
            return str(thumb_path)

    return None